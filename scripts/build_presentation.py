"""
Build the ML710 LLaVA parallelism presentation (PPTX) modeled on the
Maksym Bekuzarov 2022 example deck. ~28 slides under the 30-slide cap.

Section attribution:
  Pipeline parallelism  -> Youssef Ghallab
  ZeRO-2                -> Omar
  ZeRO-3 + FSDP         -> Rassul Magauin

Run:
    python scripts/build_presentation.py
Output:
    presentation/ML710_LLaVA_Parallelism.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


PROJECT = Path(__file__).resolve().parents[1]
PLOTS = PROJECT / "plots" / "comparison"
RUNS = PROJECT / "logs" / "runs"
OUT_DIR = PROJECT / "presentation"
OUT_DIR.mkdir(parents=True, exist_ok=True)
OUT = OUT_DIR / "ML710_LLaVA_Parallelism.pptx"

# 16:9 (matches example deck: 720x405 pt)
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]

# Colors
GREY_TXT = RGBColor(0x55, 0x55, 0x55)
DARK_TXT = RGBColor(0x10, 0x10, 0x10)
GREEN_BG = RGBColor(0xC8, 0xE6, 0xC9)
RED_BG = RGBColor(0xFF, 0xCD, 0xD2)


def add_title(slide, text, top=Inches(0.35), left=Inches(0.5),
              width=Inches(12.3), height=Inches(0.9), size=36, bold=False,
              align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = DARK_TXT
    return box


def add_bullets(slide, items, left=Inches(0.5), top=Inches(1.4),
                width=Inches(6.5), height=Inches(5.8), size=16):
    """items: list[str] or list[(text, level)] tuples."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        prefix = "    " * level + "- "
        run = p.add_run()
        run.text = prefix + text
        run.font.size = Pt(size - 2 * level)
        run.font.color.rgb = DARK_TXT if level == 0 else GREY_TXT
        p.space_after = Pt(4)
    return box


def add_image(slide, path, left, top, width=None, height=None):
    if not Path(path).exists():
        return None
    if width and height:
        return slide.shapes.add_picture(str(path), left, top, width=width, height=height)
    if width:
        return slide.shapes.add_picture(str(path), left, top, width=width)
    if height:
        return slide.shapes.add_picture(str(path), left, top, height=height)
    return slide.shapes.add_picture(str(path), left, top)


def title_slide(title, subtitle, footer=None):
    s = prs.slides.add_slide(BLANK)
    add_title(s, title, top=Inches(2.5), height=Inches(1.5),
              size=54, bold=True, align=PP_ALIGN.CENTER)
    sub = s.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.0))
    tf = sub.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = subtitle
    r.font.size = Pt(22)
    r.font.color.rgb = GREY_TXT
    if footer:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r = p2.add_run()
        r.text = footer
        r.font.size = Pt(18)
        r.font.color.rgb = GREY_TXT
    return s


def section_slide(text):
    s = prs.slides.add_slide(BLANK)
    add_title(s, text, top=Inches(3.2), height=Inches(1.2),
              size=48, bold=False, align=PP_ALIGN.CENTER)
    return s


def content_slide(title, bullets, image_path=None, image_left=Inches(7.0),
                  image_top=Inches(1.5), image_width=Inches(5.8),
                  image_height=None, bullet_width=Inches(6.3)):
    s = prs.slides.add_slide(BLANK)
    add_title(s, title)
    add_bullets(s, bullets, width=bullet_width)
    if image_path:
        add_image(s, image_path, image_left, image_top,
                  width=image_width, height=image_height)
    return s


# ============================================================
# Build the deck
# ============================================================

# 1 — Title
title_slide(
    "ML710 Project",
    "Parallelizing LLaVA Training",
    "Rassul Magauin · Omar Ahmed · Youssef Ghallab\nMBZUAI · Spring 2026",
)

# 2 — Hardware & team
content_slide(
    "Hardware & Setup",
    [
        "Cluster: MBZUAI HPC (login-student-lab.mbzu.ae)",
        ("ws-ia partition: 1× NVIDIA RTX 5000 Ada (32 GB) per node", 1),
        ("116 nodes total, 48 CPUs / 230 GB RAM each", 1),
        ("Inter-node link: ~1 Gbps Ethernet (no NVLink, no InfiniBand)", 1),
        ("gpu partition limited to 1 GPU per user (gpu-1 QOS)", 1),
        "Multi-GPU = 2 separate salloc sessions on ws-ia + torchrun --nnodes=2",
        "Software: Python 3.10, PyTorch 2.1.2, DeepSpeed 0.12.6, Transformers 4.37.2, peft 0.6.0",
        "Team: 3 students, each owning at least one parallelism strategy",
        ("Pipeline parallelism — Youssef Ghallab", 1),
        ("DeepSpeed ZeRO-2 — Omar Ahmed", 1),
        ("DeepSpeed ZeRO-3 + PyTorch FSDP — Rassul Magauin", 1),
    ],
)

# 3 — ML task: LLaVA architecture
content_slide(
    "ML task: LLaVA",
    [
        'Liu et al., "Visual Instruction Tuning", NeurIPS 2023',
        "Multimodal: image + text → text",
        "Architecture (7B params total):",
        ("CLIP ViT-L/14-336 vision encoder — frozen throughout", 1),
        ("2-layer MLP (GELU) projection — trainable, ≈80 MB", 1),
        ("Vicuna-7B LLM — frozen in Stage 1, fine-tuned in Stage 2", 1),
        "Two training stages:",
        ("Stage 1 (pre-train): MLP only on 558K image-caption pairs", 1),
        ("Stage 2 (fine-tune): MLP + LLM (or LoRA) on 665K instructions", 1),
        "Stage 2 is our main benchmark — meaningful comm volume",
        ("Stage 1 has only ~80 MB trainable; AllReduce trivial", 1),
        ("Stage 2 LoRA: ~340 MB trainable adapters, real DP cost", 1),
    ],
    image_path=None,
    bullet_width=Inches(12.3),
)

# 4 — Stage 2 dataset & training config
content_slide(
    "Stage 2 setup",
    [
        "Starting checkpoint: liuhaotian/llava-v1.5-7b",
        ("CLIP + MLP already pretrained on 558K LAION/CC/SBU pairs", 1),
        "Dataset: llava_instruct_10k.json (10K random subset, seed=42)",
        ("Source: llava_instruct_150k.json + COCO 2014 train images", 1),
        ("Subset chosen so each run finishes within ~1 hour", 1),
        "Trainable: LoRA adapters (r=128, α=256) + full MLP projection",
        ("≈4.6% of total params trainable; ~340 M of 7.4 B", 1),
        "Optimizer: AdamW; bf16; gradient checkpointing; cosine LR",
        ("LR=2e-4 (LoRA) / 5e-5 (ZeRO-3, FSDP), warmup 0.03 / 0.1", 1),
        "Effective batch sizes vary by method:",
        ("DDP / ZeRO-2: 32 (per_device=16, accum=1, world=2)", 1),
        ("Pipeline: 16 (per_device=1, accum=16, world=2)", 1),
        ("ZeRO-3 / FSDP: 8 (per_device=1, accum=4, world=2)", 1),
    ],
    bullet_width=Inches(12.3),
)

# 5 — Strategies overview
content_slide(
    "Strategies compared",
    [
        "Single-GPU LoRA (no parallelism)",
        ("Reference: how slow is no parallelism?", 1),
        "DDP (DeepSpeed ZeRO-0) — naïve baseline",
        ("Replicate model on each GPU; AllReduce gradients", 1),
        "Pipeline parallelism — model parallelism family",
        ("Split decoder layers across 2 nodes (Gloo over Ethernet)", 1),
        ("Owned by Youssef Ghallab", 1),
        "DeepSpeed ZeRO-2 — data parallelism w/ sharded optimizer + gradients",
        ("Owned by Omar Ahmed", 1),
        "DeepSpeed ZeRO-3 — data parallelism w/ sharded optimizer + grads + params",
        ("Owned by Rassul Magauin", 1),
        "PyTorch FSDP — PyTorch-native equivalent of ZeRO-3",
        ("Owned by Rassul Magauin", 1),
        "All multi-node runs use 2 ws-ia nodes × 1 GPU over Ethernet",
    ],
    bullet_width=Inches(12.3),
)

# 6 — Single-GPU baseline
content_slide(
    "Baseline: 1 GPU LoRA",
    [
        "1 RTX 5000 Ada (32 GB), no parallelism",
        ("per_device_batch=1, effective batch=1, 1 epoch over 10K", 1),
        ("10000 training steps, lr=2e-4, cosine schedule", 1),
        "Wall time: 96.0 min (5761 s)",
        "Throughput: 1.74 samples/sec",
        "Peak GPU memory: 25.7 GiB / 32 GiB (≈80%)",
        "Average GPU utilization: only 11.4%",
        ("CPU-bound: dataloader_num_workers=0 + bf16 model loading", 1),
        ("Most time in image preprocessing, not GPU compute", 1),
        "Final loss: 1.32 (avg 1.10 over the run)",
    ],
    image_path=None,
    bullet_width=Inches(12.3),
)
# We don't have a plotted curve for the 1-GPU run, just numbers.

# 7 — DDP baseline
content_slide(
    "Baseline 2: DDP (2 nodes × 1 GPU)",
    [
        "PyTorch DistributedDataParallel via torchrun --nnodes=2",
        "DeepSpeed ZeRO Stage 0 (no sharding) — pure DP",
        ("Each rank holds full model; AllReduce LoRA grads only", 1),
        "per_device_batch=16, accum=1, effective batch=32",
        ("313 training steps for 10K samples, 1 epoch", 1),
        "Wall time: 74.8 min (4491 s)",
        "Throughput: 2.22 samples/sec",
        ("vs 1-GPU @ batch 1: ~1.28× per-sample, ~10× per-step", 1),
        "Peak GPU memory: 30.7 GiB",
        "Average GPU utilization: 70.7%",
        ("Synchronization stalls on cross-node Ethernet visible", 1),
        "Final loss: 1.04 (avg 1.18)",
    ],
    image_path=RUNS / "youssef.ghallab/stage2_ddp_lora_20260420_000617/plots/loss_curve.png",
    image_width=Inches(5.8), image_top=Inches(1.7),
)

# 8 — DDP GPU plots
content_slide(
    "DDP — GPU utilization & memory",
    [
        "Both nodes track each other closely (DP is symmetric)",
        "Util oscillates ~70% during compute, drops on AllReduce",
        ("Cross-node Ethernet AllReduce dominates the gap", 1),
        "Memory steady ~30 GiB after warmup",
    ],
    image_path=RUNS / "youssef.ghallab/stage2_ddp_lora_20260420_000617/plots/gpu_utilization_curve.png",
    image_width=Inches(5.8), image_top=Inches(1.5),
)

# 9 — Section divider: Pipeline
section_slide("Pipeline Parallelism")

# 10 — Pipeline implementation
content_slide(
    "Pipeline Parallelism — implementation (Ghallab)",
    [
        "Manual 2-stage pipeline, no torchgpipe / DeepSpeed-Pipeline",
        ("Custom loop in LLaVA/llava/train/train_pipe_dist.py (522 lines)", 1),
        ("Gloo backend over Ethernet (NCCL had ordering issues with mixed dtypes)", 1),
        "Split point k = layer index in Vicuna-7B (32 decoder layers total)",
        ("Rank 0: vision tower (frozen) + projection + decoder layers 0..k-1", 1),
        ("Rank 1: decoder layers k..31 + final norm + LM head + loss", 1),
        ("Activations forward-passed and gradients backward-passed via dist send/recv", 1),
        "Microbatch + gradient accumulation: per_device=1, accum=16 → effective batch 16",
        ("num_micro_batches=1 (no inter-microbatch overlap, conservative)", 1),
        "LoRA adapters added on rank-local layers only",
    ],
    bullet_width=Inches(12.3),
)

# 11 — Pipeline split attempts
content_slide(
    "Pipeline — split balancing",
    [
        "First attempt: split=16 (50/50 by layer count)",
        ("Rank 0 OOMed in first forward — vision+embeddings+layers 0..15 = 31.4 GiB", 1),
        ("Vision tower + image features add a lot to rank 0's footprint", 1),
        "Second attempt: split=8 (shift work to rank 1)",
        ("Rank 1 OOMed — layers 8..31 + LM head = 31.0 GiB", 1),
        ("LM head + final layers heavier than expected", 1),
        "Third attempt: split=12 + microbatch=1 + accum=16",
        ("Both ranks fit; rank 1 ≈ 31.4 GiB, rank 0 ≈ 28 GiB", 1),
        ("Effective batch dropped from 16→16 (same), but per-step memory halved", 1),
        ("Completed 625/625 steps", 1),
        "Lesson: layer count is not a good proxy for memory cost.",
    ],
    bullet_width=Inches(12.3),
)

# 12 — Pipeline results
content_slide(
    "Pipeline — results",
    [
        "Wall time: 72.1 min (4325 s)",
        ("Training loop alone: 67.6 min (4058 s)", 1),
        "Throughput: 2.31 samples/sec (1.04× DDP — barely faster)",
        "Effective batch 16 (vs DDP's 32) → fewer steps per epoch",
        "Peak GPU memory: 31.4 GiB (rank 1)",
        "Average GPU utilization:",
        ("Rank 0: 33.9%  ·  Rank 1: 56.1%", 1),
        ("Big imbalance — rank 0 idle while rank 1 finishes its layers", 1),
        ("Gloo is much slower than NCCL → comm overhead amplified", 1),
        "Final loss: 1.57 (higher than DDP because of smaller effective batch)",
        "Key insight: pipeline can make a model fit, but on 1 Gbps Ethernet",
        ("activation/gradient transfers between stages dominate, killing speedup", 1),
    ],
    image_path=RUNS / "youssef.ghallab/stage2_pipe_lora_dist_20260424_153422_rank0/plots/gpu_utilization_curve.png",
    image_width=Inches(5.6), image_top=Inches(2.5),
    bullet_width=Inches(6.3),
)

# 13 — Section divider: ZeRO-2
section_slide("DeepSpeed ZeRO-2")

# 14 — ZeRO-2 explanation
content_slide(
    "ZeRO-2 — explanation (Omar)",
    [
        "Family: ZeRO-DP (Zero Redundancy Optimizer Data Parallel)",
        "Each rank still processes a different microbatch (data parallelism)",
        "What's sharded across the 2 ranks:",
        ("Optimizer states (AdamW m, v) — sharded ✓", 1),
        ("Gradients (after backward) — sharded ✓", 1),
        ("Parameters (forward + backward weights) — replicated ✗", 1),
        "Communication pattern per step:",
        ("Backward: ReduceScatter gradients to owner-rank shards", 1),
        ("Optimizer step: each rank updates its shard of params", 1),
        ("AllGather updated params back to full set on each rank", 1),
        "Memory savings vs DDP:",
        ("Optimizer state for LoRA only ≈ a few hundred MB → little to shard", 1),
        ("Real memory savings come from offloaded full-precision adam states", 1),
        "Config: configs/zero2.json, launcher: scripts/run_stage2_zero2_lora.sh",
    ],
    bullet_width=Inches(12.3),
)

# 15 — ZeRO-2 results
content_slide(
    "ZeRO-2 — results",
    [
        "Run: stage2_zero2_lora_20260424_134311 (effective batch 32)",
        "Wall time: 58.4 min (3503 s)",
        ("vs DDP 74.8 min — saved 16.4 min, −22% wall-clock", 1),
        "Throughput: 2.85 samples/sec — 1.28× DDP",
        "Peak GPU memory: 27.6 GiB (vs DDP 30.7 GiB) — saved 3.1 GiB",
        "Average GPU utilization: 88.5% (vs DDP 70.7%)",
        ("DeepSpeed's bucketed async comm overlaps better than DDP", 1),
        ("Better overlap fills the AllReduce-stall gap DDP suffered", 1),
        "Final loss: 1.11 (avg 1.22) — comparable to DDP",
        "Why only 1.28× and not 2×?",
        ("LoRA's optimizer state is already tiny — sharding gains are small", 1),
        ("Most of the win comes from DeepSpeed's overlap, not from sharding", 1),
        "ZeRO-2 = the only method that beats DDP on this workload.",
    ],
    image_path=RUNS / "youssef.ghallab/stage2_zero2_lora_20260424_134311/plots/gpu_utilization_curve.png",
    image_width=Inches(5.6), image_top=Inches(2.0),
    bullet_width=Inches(6.3),
)

# 16 — Section divider: ZeRO-3
section_slide("DeepSpeed ZeRO-3")

# 17 — ZeRO-3 explanation
content_slide(
    "ZeRO-3 — explanation (Rassul)",
    [
        "Family: ZeRO-DP — same data-parallel semantics as ZeRO-2",
        "What's sharded:",
        ("Optimizer states ✓", 1),
        ("Gradients ✓", 1),
        ("Parameters ✓ ← new vs ZeRO-2", 1),
        "Communication pattern per step (and per layer!):",
        ("Forward layer N: AllGather params for layer N → compute → drop", 1),
        ("Backward layer N: AllGather params again → compute grads → drop", 1),
        ("ReduceScatter gradients at end of backward", 1),
        "On 2 ranks the model is split in half (~7 GB sharded each)",
        ("But every layer requires re-materializing the full ≈430 MB to each rank", 1),
        "Math for our workload:",
        ("32 transformer layers × 430 MB × 2 (fwd+bwd) ≈ 27.5 GB transferred per step", 1),
        ("Over 1 Gbps Ethernet ≈ 4 minutes of comm per step alone", 1),
        "Config: configs/zero3.json + zero3_bounded_async.json (50 MB buckets)",
    ],
    bullet_width=Inches(12.3),
)

# 18 — ZeRO-3 results
content_slide(
    "ZeRO-3 — results",
    [
        "Run: stage2_zero3_lora_20260420_041514/041515 (2 ranks)",
        "Hit SLURM time limit before completion",
        ("Reached 292/1250 steps after 3 h 52 min", 1),
        ("Each step ≈ 48 s; projected full wall time ≈ 16.5 h", 1),
        "Throughput: 0.17 samples/sec — 0.08× DDP (≈13× slower)",
        "Effective batch 8 (per_device=1, accum=4)",
        ("Smaller batch chosen because GPU mem accounting is more conservative under ZeRO-3", 1),
        "Peak GPU memory: 22.0 GiB",
        ("Lower than DDP (30.7) — sharding does save memory", 1),
        ("But the memory wasn't the bottleneck; comm was", 1),
        "Average GPU utilization: 99.2%",
        ("Counter-intuitive: 99% util but slowest — DeepSpeed pipelines overlap", 1),
        ("comm with compute, so 'utilized' time includes waiting on AllGather", 1),
        "Last logged loss: 1.31 — convergence direction OK, just too slow",
        "Diagnosis: per-layer cross-node AllGather is the bottleneck.",
    ],
    image_path=None,
    bullet_width=Inches(12.3),
)

# 19 — Section divider: FSDP
section_slide("PyTorch FSDP")

# 20 — FSDP explanation
content_slide(
    "PyTorch FSDP — explanation (Rassul)",
    [
        'PyTorch-native equivalent of DeepSpeed ZeRO-3',
        ("Same algorithm: shard params + grads + optim across DP ranks", 1),
        ("Different implementation: PyTorch core, not DeepSpeed", 1),
        "Different design choices vs ZeRO-3:",
        ("FSDP wraps each transformer block as its own FSDP unit (FlatParameter)", 1),
        ("AllGather happens at unit boundary, not per-tensor", 1),
        ("Mixed precision policy is explicit (we used bf16)", 1),
        "Config: configs/fsdp.json",
        ("auto_wrap_policy=TRANSFORMER_BASED_WRAP, wrap LlamaDecoderLayer", 1),
        ("use_orig_params=True (required for LoRA non-uniform requires_grad)", 1),
        ("cpu_ram_efficient_loading=true (only rank 0 holds full CPU state at load)", 1),
        "Two LLaVA bugs we had to patch:",
        ("vision_tower[0] fails when loaded from a vendored checkpoint (not a list)", 1),
        ("LoRA adapters created in fp32 after model.to(bf16) → cast LoRA to bf16 too", 1),
        "Launcher: scripts/run_stage2_fsdp_lora.sh (mirrors run_stage2_zero3_lora.sh)",
    ],
    bullet_width=Inches(12.3),
)

# 21 — FSDP results
content_slide(
    "FSDP — results",
    [
        "Run: stage2_fsdp_lora_20260426_143948",
        "Stopped early at step 28/1250 (matches the ZeRO-3 timeout pattern)",
        "Wall time observed: 39 min (2350 s); per-step ≈ 80 s",
        ("Projected full wall time ≈ 27.9 h — even slower than ZeRO-3", 1),
        "Throughput: 0.095 samples/sec — 0.04× DDP, ≈22× slower",
        ("≈1.7× slower than ZeRO-3 on the same workload, same hardware", 1),
        "Peak GPU memory: 32.3 GiB — basically saturating the 32 GiB card",
        ("FSDP's FlatParameter contiguous tensor is larger than DeepSpeed's bucket scheme", 1),
        "Average GPU utilization: 97.5%",
        "Why even slower than ZeRO-3?",
        ("DeepSpeed ZeRO-3 with bounded async config (50 MB buckets) overlaps better", 1),
        ("FSDP defaults: no forward_prefetch, no bucket tuning → strict serial gather", 1),
        ("FSDP gathers per transformer block (~430 MB), DeepSpeed by 50 MB buckets", 1),
        "Same fundamental wall: 1 Gbps Ethernet AllGather of bf16 weights every layer.",
    ],
    image_path=RUNS / "rassul.magauin/stage2_fsdp_lora_20260426_143948/plots/gpu_utilization_curve.png",
    image_width=Inches(5.6), image_top=Inches(1.5),
    bullet_width=Inches(6.3),
)

# 22 — Section divider: Comparison
section_slide("Cross-method comparison")

# 23 — Throughput
content_slide(
    "Throughput — samples/sec (log scale)",
    [
        "ZeRO-2: 2.85 sps — only method beating DDP",
        "DDP:    2.22 sps — naïve baseline",
        "Pipeline: 2.31 sps — barely above DDP",
        "1-GPU LoRA: 1.74 sps — surprisingly close to multi-GPU at batch 1",
        "ZeRO-3: 0.17 sps — 13× slower (timed out)",
        "FSDP:   0.095 sps — 22× slower (stopped early)",
    ],
    image_path=PLOTS / "1_throughput.png",
    image_width=Inches(7.0), image_top=Inches(1.4),
    bullet_width=Inches(5.0),
)

# 24 — Speedup vs DDP
content_slide(
    "Speedup vs DDP baseline",
    [
        "DDP = 1.00× by definition",
        "ZeRO-2: 1.28× — only winner",
        "Pipeline: 1.04× — basically tied with DDP",
        "1-GPU LoRA: 0.78× — naïve scaling reference",
        "ZeRO-3: 0.08× — 13× regression",
        "FSDP: 0.04× — 22× regression",
        "Pattern: DP variants that don't shard params win;",
        ("methods that shard params lose to Ethernet overhead", 1),
    ],
    image_path=PLOTS / "3_speedup_vs_ddp.png",
    image_width=Inches(7.0), image_top=Inches(1.4),
    bullet_width=Inches(5.0),
)

# 25 — Memory
content_slide(
    "Peak GPU memory",
    [
        "32 GiB hardware limit (RTX 5000 Ada)",
        "ZeRO-2: 27.6 GiB — best of the multi-GPU methods",
        "DDP: 30.7 GiB",
        "Pipeline (k=12): 31.4 GiB — heavier rank near the limit",
        "FSDP: 32.3 GiB — saturated; FlatParameter overhead",
        "ZeRO-3: 22.0 GiB — actual sharding visible (but at huge speed cost)",
        "1-GPU LoRA: 25.7 GiB — fits comfortably without parallelism",
    ],
    image_path=PLOTS / "4_gpu_memory.png",
    image_width=Inches(7.0), image_top=Inches(1.4),
    bullet_width=Inches(5.0),
)

# 26 — Goodput
content_slide(
    "Goodput — useful learning per cost",
    [
        "Loss reduction per wall-clock second:",
        ("ZeRO-2: 0.000372  ← best", 1),
        ("DDP: 0.000313", 1),
        ("ZeRO-3: 0.000294 (still positive — converging, just slowly)", 1),
        ("Pipeline: 0.000239", 1),
        ("FSDP: −0.000056 (only 28 steps; started from already-tuned ckpt)", 1),
        "Loss reduction per training sample (statistical efficiency):",
        ("DDP: 0.0045 ← best per-sample (effective batch 32)", 1),
        ("ZeRO-2: 0.0042 (also batch 32)", 1),
        ("ZeRO-3 / Pipeline: ~0.0017 (smaller effective batch hurts)", 1),
        "Wall-clock goodput tracks what slides actually want.",
    ],
    image_path=PLOTS / "9_goodput.png",
    image_width=Inches(7.0), image_top=Inches(1.6),
    bullet_width=Inches(5.0),
)

# 27 — Loss curves
content_slide(
    "Loss vs step (smoothed)",
    [
        "1-GPU LoRA: gray, 10000 steps, slowly converging to ~1.0",
        "DDP / ZeRO-2: ~313 steps each, converge to ~1.04 / 1.11",
        "Pipeline: 625 steps, slightly higher final loss (smaller batch)",
        "ZeRO-3: dashed red, only 292 steps reached, last loss 1.31",
        "FSDP: dashed purple, 28 steps, started from pretrained ckpt",
        "All methods that completed reach comparable loss",
        ("Statistical efficiency is similar; differences are wall-clock", 1),
    ],
    image_path=PLOTS / "6_loss_vs_step.png",
    image_width=Inches(7.0), image_top=Inches(1.4),
    bullet_width=Inches(5.0),
)

# 28 — Key findings
content_slide(
    "Key findings",
    [
        "On 1 Gbps Ethernet with a 7B-LoRA workload:",
        "1. ZeRO-2 wins (1.28×) — sharding optimizer + grads, not params, is the sweet spot.",
        ("LoRA already shrinks the gradient AllReduce → DDP's overhead is the bottleneck.", 1),
        ("DeepSpeed's bucketed async overlap fills the gap better than vanilla DDP.", 1),
        "2. Pipeline parallelism (1.04×) makes the model fit but doesn't speed things up.",
        ("Cross-node Gloo activation transfers are slower than DDP's AllReduce.", 1),
        ("Layer count is not a good proxy for memory load — split balancing is empirical.", 1),
        "3. ZeRO-3 / FSDP collapse on Ethernet (0.04×–0.08×).",
        ("Per-layer AllGather of full weights every step is intractable over 1 Gbps.", 1),
        ("Implementation matters: DeepSpeed ZeRO-3 is ~1.7× faster than PyTorch FSDP", 1),
        ("because of better bucket tuning and async overlap defaults.", 1),
        "4. ZeRO-3 / FSDP would scale on InfiniBand or NVLink — wrong tool for our hardware.",
        "5. LoRA training on a 7B model fits on a single RTX 5000 Ada — choose DP variants.",
    ],
    bullet_width=Inches(12.3),
)

# 29 — Lessons learned
content_slide(
    "Lessons learned",
    [
        "Match the parallelism strategy to the bottleneck",
        ("Optimizer-state-bound? → ZeRO-2", 1),
        ("Gradient-comm-bound? → DDP with good overlap is hard to beat", 1),
        ("Activation-memory-bound (long sequence)? → Pipeline / TP, not relevant for us", 1),
        ("Param-memory-bound (model > GPU)? → ZeRO-3 / FSDP, but only with fast interconnect", 1),
        "Implementation details actually matter when comm-bound",
        ("Bucket sizes, async overlap, prefetch policy can move 1.7× of perf", 1),
        ("DeepSpeed > FSDP defaults on slow Ethernet", 1),
        "Memory headroom is not a free win",
        ("ZeRO-3 saved 9 GiB but cost us 13× wall-clock — only useful if model doesn't fit", 1),
        "LoRA breaks naïve assumptions",
        ("Trainable params drop from 7 B to 340 M → optimizer/grad sharding is small wins", 1),
        ("Mixed dtypes (bf16 base + fp32 LoRA) need explicit handling under FSDP", 1),
        "Cross-node infrastructure dominates",
        ("1 Gbps Ethernet is the true bottleneck, not DeepSpeed vs FSDP vs Pipeline", 1),
    ],
    bullet_width=Inches(12.3),
)

# 30 — Thanks
title_slide(
    "Thank you",
    "Questions?",
    "Repo: github.com/rassulmagauin/ml710\nReport: ZERO2_RESULTS.md · results_explanation.md",
)


prs.save(OUT)
print(f"Wrote {OUT}")
print(f"Slides: {len(prs.slides)}")
